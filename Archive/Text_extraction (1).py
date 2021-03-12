 #!/usr/bin/env python
# coding: utf-8

'''

BEFORE USE
Search for 'UPDATE THE' and update the path, or folder structure you are using.
There are 4 sections to update all found in the last 15-20% of the code.

This is a text extract facility from the following file formats:

.PDF
.MSG
.DOCX
.DOC
.PPT
.PPTX
.XLS
.XLSX
.TXT

The facility will open the file on your laptop, process it and close it.
Note that PDFs and PowerPoint Files use objects and so need to be converted and read, so not as straight forward as thought.

'''

import pickle
import os
import re
import pandas as pd
import csv
import time
from progressbar import ProgressBar
import io
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from pdfminer.converter import TextConverter
from pdfminer.pdfinterp import PDFPageInterpreter
from pdfminer.pdfinterp import PDFResourceManager
from pdfminer.pdfpage import PDFPage
from pdfminer.layout import LAParams
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LTTextBoxHorizontal, LTTextBox, LTTextLine, LTChar, LTTextLineHorizontal, LTRect
import win32com
from win32com.client import Dispatch
import docx
from docx.document import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
import extract_msg
from pptx import Presentation
import xlrd
import sys
import glob
import matplotlib.pyplot as plt
from matplotlib import patches
from datetime import date

#  Define text extraction functions
# TXT files


def extract_text_txt(file_path, text_string = None):
    if text_string is None:
        text_string = ''
    try:
        with open(file_path) as fh:
            text = fh.read().replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ')
            text_string += text + ' '
        return(text_string)
    except Exception as e:
        print(f'{file_path} could not be parsed')

## The following pdf parsing functions are archived as a more complete solution was found

#def page_text_extract(path_to_pdf):
#    '''
#    Extracts the text for each page in a pdf document
#    '''
#    with open(path_to_pdf, 'rb') as open_path:
#        for page in PDFPage.get_pages(open_path, 
#                                      caching=True,
#                                      check_extractable=True):
#            resource_manager = PDFResourceManager()
#            fake_file_handle = io.StringIO()
#            converter = TextConverter(resource_manager, fake_file_handle)
#            page_interpreter = PDFPageInterpreter(resource_manager, converter)
#            page_interpreter.process_page(page)
#            text = fake_file_handle.getvalue()
#            yield text
#    
#            # close open handles
#            converter.close()
#            fake_file_handle.close()
#    
#def extract_text_pdf(path_to_pdf, text_str = None):
#    '''
#    Extracts text for an entire pdf document
#    '''
#    if text_str is None:
#        text_str = ''
#    try:
#        for page in page_text_extract(path_to_pdf):
#            text_str += page + ' '
#        return(text_str)
#    except:
#        print(f'{path_to_pdf} could not be parsed')


#PDF files

def extract_layout_by_page(pdf_path, layouts = None):
    """
    Extracts LTPage objects from a pdf file.
    
    slightly modified from
    https://euske.github.io/pdfminer/programming.html
    """
    laparams = LAParams()
    with open(pdf_path, 'rb') as fp:
        parser = PDFParser(fp)
        document = PDFDocument(parser)
        rsrcmgr = PDFResourceManager()
        device = PDFPageAggregator(rsrcmgr, laparams=laparams)
        interpreter = PDFPageInterpreter(rsrcmgr, device)

        if layouts is None:
            layouts = []
            
        for page in PDFPage.create_pages(document):
            interpreter.process_page(page)
            layouts.append(device.get_result())

    return layouts

def flatten(lst):
    """Flattens a list of lists"""
    return [subelem for elem in lst for subelem in elem]


def extract_characters(element, text_elements = None):
    """
    Recursively extracts individual characters from 
    text elements. 
    """
    if text_elements is None:
        TEXT_ELEMENTS = [
                        LTChar,
                        LTTextBox,
                        LTTextBoxHorizontal,
                        LTTextLine,
                        LTTextLineHorizontal]
        
    if isinstance(element, LTChar):
        return [element]

    if any(isinstance(element, i) for i in TEXT_ELEMENTS):
        return flatten([extract_characters(e) for e in element])

    if isinstance(element, list):
        return flatten([extract_characters(l) for l in element])

    return []

# The following two functions are used to plot a visual representation of the pdf parsing, they are not 
# used in the extraction

def draw_rect_bbox(tup, ax, color):
    """
    Draws an unfilled rectable onto ax.
    """
    if isinstance(tup, tuple):
        x0,y0,x1,y1 = tup[0],tup[1],tup[2],tup[3]
    ax.add_patch( 
        patches.Rectangle(
            (x0, y0),
            x1 - x0,
            y1 - y0,
            fill=False,
            color=color
        )    
    )
    
def draw_rect(rect, ax, color="black"):
    draw_rect_bbox(rect.bbox, ax, color)
    
def new_pdf_extraction(file_path, text_string = None,text_string_main = None, texts = None, rects = None, coord_list = None):
    try:
        page_layouts = extract_layout_by_page(file_path)
        
        if text_string_main is None:
            text_string_main = ''
        for current_page in page_layouts:
            texts = None
            rects = None
            coord_list = None
            text_string = None
            if text_string is None:
                text_string = ''
            if texts is None:
                texts = []
            if rects is None:
                rects = []
            if coord_list is None:
                coord_list = []
            # seperate text and rectangle elements
            for e in current_page:
                if isinstance(e, LTTextBoxHorizontal):
                    texts.append(e)
                elif isinstance(e, LTRect):
                    rects.append(e)
        
            # sort them into 
            current_characters = extract_characters(texts)

# The following code is used for visual representation of the extraction
            
#            xmin, ymin, xmax, ymax = current_page.bbox
#            size = 6
#            
#            fig, ax = plt.subplots(figsize = (size, size * (ymax/xmax)))
#            
#            for rect in rects:
#                draw_rect(rect, ax)
#                
#            for c in current_characters:
#                draw_rect(c, ax, "red")
#                
#            
#            plt.xlim(xmin, xmax)
#            plt.ylim(ymin, ymax)
#            plt.show()
            
            for c in current_characters:
                x_coord = c.matrix[4]
                y_coord = round(c.matrix[5],0)
                text = c.get_text()
                char_list = [x_coord,y_coord,text]
                coord_list.append(char_list)
                
            sorted_list = sorted(coord_list, key = lambda x : (x[1],-x[0]), reverse = True)
            
            if sorted_list:
                init_y = sorted_list[0][1]
                init_x = sorted_list[0][0]
                i = 0
                for coord in sorted_list:
                    if coord[1] == init_y and (coord[0]-init_x) < 20:
                        text_string += coord[2].replace('\xa0',' ')
                        init_y = coord[1]
                        init_x = coord[0] 
                    else:
                        spaced_text = ' ' + coord[2].replace('\xa0',' ')
                        text_string+=spaced_text
                        init_y = coord[1]
                        init_x = coord[0]
            
                text_string_main+=text_string
        if not text_string_main:
            print(f'{file_path} was parsed, however, contains no text. This could be due to the file being comprised of only images')
        return(text_string_main)
    
    except Exception as e:
        print(f'{file_path} could not be parsed')
        
# DOC & DOCX files (with table text extraction)

def doc2docx(path):
    '''
    Convert doc file to docx if necessary
    '''
    word = Dispatch('word.application')
    word.DisplayAlerts = 0
    word.visible = 0
    doc = word.Documents.Open(path)
    doc.SaveAs(path+"x", 12)
    doc.Close()
    word.Quit()
    
    #os.remove(path)
    
def iter_block_items(parent):
    """
    Yield each paragraph and table child within *parent*, in document
    order. Each returned value is an instance of either Table or
    Paragraph.
    """
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    elif isinstance(parent, docx._Cell):
        parent_elm = parent._tc
    else:
        raise ValueError(f"Docx paring issue, please check document {parent}")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)
            
def extract_text_doc_x(file_path, text_string = None):
    
    if text_string is None:
        text_string = ''
    doc_re = re.compile('.doc$|.DOC$')
    
    try:
        if doc_re.search(file_path):
            docx_version = doc2docx(file_path)
            file_path = file_path +'x'
                
        document = docx.Document(file_path)
        for block_item in iter_block_items(document):
            if isinstance(block_item, Paragraph):
                paratext = block_item.text.replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ') + ' '
                text_string += paratext
            elif isinstance(block_item, Table):
                for row in block_item.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            table_text = paragraph.text.replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ') + ' '
                            text_string += table_text
            else:
                pass
        return(text_string)
    
    except Exception as e:
        print(f'{file_path} could not be parsed')


## The following docx parsing functions are archived as a more complete solution was found, which included table parsing
    
# def extract_docx_text(doc_or_docx_path, text_string = None):
#     '''
#     Print text from docx file
#     '''
#     if text_string is None:
#         text_string = ''
        
#     doc_re = re.compile('.doc$|.DOC$')
#     if doc_re.search(doc_or_docx_path):
#         docx_version = doc2docx(doc_or_docx_path)
#         doc_or_docx_path = doc_or_docx_path+'x'
        
#     doc = docx.Document(doc_or_docx_path)
#     for p in doc.paragraphs:
#         text = p.text
#         text_string += text
#     return(text_string)

# MSG files
        
def extract_text_msg(file_path, text_string = None):
    if text_string is None:
        text_string = ''
    try:
        msg = extract_msg.Message(file_path)
        msg_sender = msg.sender
        msg_date = msg.date
        msg_subj = msg.subject
        msg_message = msg.body
        text_string+=msg_sender.replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ')
        text_string+=msg_date.replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ')
        text_string+=msg_subj.replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ')
        text_string+=msg_message.replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ')
        return(text_string)
    except Exception as e:
        print(f'{file_path} could not be parsed')
    

#xlxs & xls files 

def extract_text_xls_x(file_path, text_string = None):
    if text_string is None:
        text_string = ''
    try:    
        book = xlrd.open_workbook(file_path)   
        for sheet in book.sheets():
            for i in range(sheet.nrows):
                row_values = sheet.row_values(i)
                for val in row_values:
                    if val:
                        text_string += str(val).replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ')
                        text_string += ' '
    
        return(text_string)
    except Exception as e:
        print(f'{file_path} could not be parsed')

       
# PPTX files

def extract_text_pptx(file_path, text_string = None):
    ppt_re = re.compile('.ppt$|.PPT$')
    if text_string is None:
        text_string = ''
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = shape.text.replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ') + ' '
                    text_string += shape_text
        return(text_string)
    except:
        print(f'{file_path} could not be parsed')

# PPT files

def convert(filename, formatType = 32):  
    powerpoint = Dispatch("Powerpoint.Application")
    powerpoint.Visible = 1    
    newname = os.path.splitext(filename)[0] + ".pdf"  
    deck = powerpoint.Presentations.Open(filename)  
    deck.SaveAs(newname, formatType)  
    deck.Close()  
    powerpoint.Quit()
    return(newname)

def extract_pdf_ppt_document(document, text_string = None):
    if text_string is None:
        text_string = ''
    rsrcmgr = PDFResourceManager()
    laparams = LAParams()
    device = PDFPageAggregator(rsrcmgr, laparams=laparams)
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    with open(document, 'rb') as document:
        for page in PDFPage.get_pages(document):
                interpreter.process_page(page)
                layout = device.get_result()
                for element in layout:
                    if isinstance(element, LTTextBoxHorizontal):
                        element_text = element.get_text().replace('\n', ' ').replace('\r', ' ').replace('\xa0', ' ').replace('\t',' ').replace('•', ' ') + ' '
                        text_string += element_text
    return(text_string)
    
def extract_text_ppt(file_path):
    try:
        path_string = convert(file_path)
        ppt_pdf_path = r"{}".format(path_string)
        text_string = extract_pdf_ppt_document(ppt_pdf_path)
        return(text_string)
    except Exception as e:
        print(f'{file_path} could not be parsed')

#
# - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - USER INPUT REQUIRED
#  __  UPDATE THE folder names before running,
#  __  OR, set the root varaible as your working directory
#


user = input("Please enter your pc user name: ")
# describe file path to your chosen folder where the documents are stored
root = 'C:\\Users\\' + user + '\\<enter foler namne>\\<enter folder namne>\\<enter folder namne>'
dirlist = [ item for item in os.listdir(root) if os.path.isdir(os.path.join(root, item)) ]

print('No. of documents: {}'.format(len(dirlist)))

file_list = []

for dir in dirlist:
    curr_files = [ item for item in os.listdir(root + '\\' + dir) ]
    file_list.append(curr_files)
    
files = list(zip(dirlist, file_list))

with pd.option_context('display.max_rows', None, 'display.max_columns', None, 'display.max_colwidth', 0):
    df = pd.DataFrame(files, columns = ['Documents','Files'])
    new_df = df.style.set_properties(**{'text-align': 'left'})

#     display(new_df)

i=0
for inc in files:
    for file in inc[1]:
        i+=1
print(f'No. of documents: {i}')



pbar = ProgressBar()
pdf_re = re.compile('.pdf$|.PDF$')
msg_re = re.compile('.msg$|.MSG$')
docx_re = re.compile('.docx$|.DOCX$|.DOCx$')
doc_re = re.compile('.doc$|.DOC$')
ppt_re = re.compile('.ppt$|.PPT$')
pptx_re = re.compile('.pptx$|.PPTX$')
xls_re = re.compile('.xls$|.XLS$')
xlsx_re = re.compile('.xlsx$|.XLSX$')
txt_re = re.compile('.txt$|.TXT$')

text_df = pd.DataFrame(columns=['Document','File extension','File Name','File Text'])

#
#  UPDATE THE path VARIABLE BELOW
#

for inc in pbar(files):
    for file in inc[1]:
        path =r'C:\Users\{user}\<enter foler namne>\<enter folder namne>\<enter folder namne>\{inc[0]}\{file}'.format(user=user, inc = inc, file=file)
        
        if pdf_re.search(file):
            extracted_pdf_text = new_pdf_extraction(path)
            text_df = text_df.append({'Doc':inc[0], 'File extension':'PDF','File Name':file,'File Text':extracted_pdf_text}, ignore_index=True)
        elif msg_re.search(file):
            extracted_msg_text = extract_text_msg(path)
            text_df = text_df.append({'Doc':inc[0], 'File extension':'MSG','File Name':file,'File Text':extracted_msg_text}, ignore_index=True)
            
        elif docx_re.search(file) or doc_re.search(file):
            extracted_docx_text = extract_text_doc_x(path)
            text_df = text_df.append({'Doc':inc[0], 'File extension':'DOC/DOCX','File Name':file,'File Text':extracted_docx_text}, ignore_index=True)
            
        elif pptx_re.search(file):
            extracted_pptx_text = extract_text_pptx(path)
            text_df = text_df.append({'Doc':inc[0], 'File extension':'PPTX','File Name':file,'File Text':extracted_pptx_text}, ignore_index=True)
            
        elif xlsx_re.search(file) or xls_re.search(file):
            extracted_xls_x_text = extract_text_xls_x(path)
            text_df = text_df.append({'Doc':inc[0], 'File extension':'XLS/XLSX','File Name':file,'File Text':extracted_xls_x_text}, ignore_index=True)
            
        elif txt_re.search(file):
            extracted_txt_text = extract_text_txt(path)
            text_df = text_df.append({'Doc':inc[0], 'File extension':'TXT','File Name':file,'File Text':extracted_txt_text}, ignore_index=True)
            
        elif ppt_re.search(file):
            extracted_ppt_text = extract_text_ppt(path)
            text_df = text_df.append({'Incident':inc[0], 'File extension':'PPT','File Name':file,'File Text':extracted_ppt_text}, ignore_index=True)              
        else:
            print('Unable to extract {} from doc {}'.format(file, inc[0]))


no_dupes_text_df = text_df.drop_duplicates(['Doc','File Text'])
mask = (no_dupes_text_df['File Text'].str.len() <= 10)
null_text_df = no_dupes_text_df.loc[mask]
join_df = no_dupes_text_df.merge(null_text_df, on=['Doc','File Name'], how='left', indicator=True)
final_text_df = join_df[join_df['_merge'] == 'left_only']
final_text_df.drop(columns=['File extension_y','File Text_y','_merge'], inplace=True)
final_text_df.dropna(subset = ['File Text_x'], inplace=True)

#
#  UPDATE THE path BELOW - where you want the output file written
#

today = date.today()
final_text_df.to_csv(r'C:\Users\{user}\<enter foler namne>\<enter folder namne>\<enter folder namne>\Text_extraction_{today}.csv'.format(user = user, today = today),index=False)

master_list = []

group_incidents_df = final_text_df.groupby('Incident')

for k, v in group_incidents_df:
    inc_list = []
    for t in v['File Text_x']:
        inc_list.append(t)
    master_list.append( tuple([k,' '.join(inc_list)]) )

final_grouped_df = pd.DataFrame(master_list, columns=['Doc','Text'])

#
#  UPDATE THE inc_extract_path VARIABLE BELOW
#

for index, row in final_grouped_df.iterrows():
    inc_no = row[0]
    inc_text = row[1]
    inc_extract_path = r'C:\Users\{user}\<enter foler namne>\<enter folder namne>\<enter folder namne>\{inc_no}.txt'.format(user=user, inc_no = inc_no)
    if os.path.isfile(inc_extract_path):
        with open(inc_extract_path, "a", encoding='utf-8') as txt_file:
            txt_file.write(" " + inc_text)
            txt_file.close()
    else:
        with open(inc_extract_path, "w", encoding='utf-8') as txt_file:
            txt_file.write(inc_text)
            txt_file.close()
        

#text_df[text_df['File extension'] == 'DOC/DOCX'].drop_duplicates(['Incident','File extension','File Text'])
# these are the dupes
#dupes_df = text_df[text_df.duplicated(['Incident','File extension','File Text'])]

